from pathlib import Path
try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from ModuleFolders.Infrastructure.Cache.CacheFile import CacheFile
from ModuleFolders.Domain.FileOutputer.BaseWriter import (
    BaseTranslatedWriter,
    OutputConfig,
    PreWriteMetadata
)

class PptxWriter(BaseTranslatedWriter):
    def __init__(self, output_config: OutputConfig):
        super().__init__(output_config)

    @classmethod
    def get_project_type(cls):
        return "Pptx"

    def on_write_translated(
        self, translation_file_path: Path, cache_file: CacheFile,
        pre_write_metadata: PreWriteMetadata,
        source_file_path: Path = None,
    ):
        if not source_file_path or not source_file_path.exists():
            return

        try:
            # 加载源文件作为模板
            prs = Presentation(source_file_path)
            
            # 建立映射以便快速写入： (slide_idx, shape_id, para_idx) -> translated_text
            translation_map = {}
            for item in cache_file.items:
                if item.final_text:
                    key = (
                        item.get_extra("slide_idx"),
                        item.get_extra("shape_id"),
                        item.get_extra("para_idx")
                    )
                    translation_map[key] = item.final_text

            # 再次遍历并替换文本
            for slide_idx, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        key = (slide_idx, shape.shape_id, para_idx)
                        if key in translation_map:
                            # 直接替换段落文本会丢失部分富文本格式（加粗/变色），但这是最稳定的翻译替换方式
                            paragraph.text = translation_map[key]

            # 保存文件
            prs.save(translation_file_path)
            
        except Exception as e:
            print(f"Error writing PPTX: {e}")