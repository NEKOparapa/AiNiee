from pathlib import Path
from pptx import Presentation


from ModuleFolders.Infrastructure.Cache.CacheFile import CacheFile
from ModuleFolders.Infrastructure.Cache.CacheItem import CacheItem, TranslationStatus
from ModuleFolders.Domain.FileReader.BaseReader import (
    BaseSourceReader,
    InputConfig,
    PreReadMetadata
)

class PptxReader(BaseSourceReader):
    def __init__(self, input_config: InputConfig):
        super().__init__(input_config)
        if Presentation is None:
            raise ImportError("Please install python-pptx library: pip install python-pptx")

    @classmethod
    def get_project_type(cls):
        return "Pptx"

    @property
    def support_file(self):
        return "pptx"

    def on_read_source(self, file_path: Path, pre_read_metadata: PreReadMetadata) -> CacheFile:
        items = []
        try:
            prs = Presentation(file_path)
            
            # 遍历幻灯片
            for slide_idx, slide in enumerate(prs.slides):
                # 遍历形状
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    
                    # 遍历段落
                    for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                        text = paragraph.text
                        if text and text.strip():
                            item = CacheItem(
                                source_text=text,
                                translation_status=TranslationStatus.UNTRANSLATED,
                                extra={
                                    "slide_idx": slide_idx,
                                    "shape_id": shape.shape_id,
                                    "para_idx": para_idx
                                }
                            )
                            items.append(item)
                            
        except Exception as e:
            print(f"Error reading PPTX {file_path}: {e}")
            return None

        return CacheFile(items=items)